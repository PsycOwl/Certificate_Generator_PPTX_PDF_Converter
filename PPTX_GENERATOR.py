from pptx import Presentation
import xlrd

loc = ("Name_list.xlsx") # <--- The Excel file of names
wb = xlrd.open_workbook(loc)
sheet = wb.sheet_by_index(0)

sheet.cell_value(0, 0) # Header



def XLStoPPTX(DocumentID, Name):
    prs = Presentation('Certificate_Template.pptx')
    for shape in prs.slides[0].shapes:
        if (shape.has_text_frame):
            if(shape.text_frame.text == 'Name'):
                shape.text_frame.paragraphs[0].runs[0].text = Name
                print(Name)
            if(shape.text_frame.text == 'DocumentID'):
                shape.text_frame.paragraphs[0].runs[0].text = DocumentID
                print(DocumentID)
    prs.save('GENERATED_PPTX/'+ DocumentID + '_' + Name + '.pptx')

for i in range(sheet.nrows):
    if (i == 0):
        continue # Ignore Header of the excel file
    DocumentID = "{:02d}".format(i)
    Name = str(sheet.cell_value(i, 0)).strip()
    XLStoPPTX(DocumentID, Name)